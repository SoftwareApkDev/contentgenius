import os
from pptx import Presentation

from llm.prompt_builder import generate_response


def extract_text_from_ppt(ppt_path):
    """
    Extracts all text content from a PowerPoint (.pptx) file.
    """
    if not os.path.exists(ppt_path):
        raise FileNotFoundError(f"File not found: {ppt_path}")

    prs = Presentation(ppt_path)
    text_chunks = []

    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text.append(shape.text.strip())
        slide_content = "\n".join(slide_text).strip()
        if slide_content:
            text_chunks.append(f"Slide {i+1}:\n{slide_content}")

    return "\n\n".join(text_chunks)


def build_prompt(slide_text: str, question: str):
    """
    Constructs a prompt for the LLM using the extracted PowerPoint content.
    """
    return f"{question}\n\n---\n\nPowerPoint content:\n\n{slide_text}"


def analyze(ppt_path: str, question: str):
    """
    Main analysis function for PowerPoint presentations.
    """
    print(f"📽️ Reading PowerPoint: {ppt_path}")
    text = extract_text_from_ppt(ppt_path)

    if not text.strip():
        return "❌ No readable text found in the presentation."

    prompt = build_prompt(text[:6000], question)  # Trim if very long
    print("🧠 Sending to LLM...")
    response = generate_response(prompt)
    return response
